"""Generate a project PPTX using a reference template PPTX.

This script loads the user's reference PPTX (sampleppt (3).pptx), replaces text
in the first 7 slides with the drought prediction project content, and saves a
new PPTX in the project root.

Notes:
- It preserves the template's theme and layout as much as possible.
- It only edits text placeholders (title/body). Other design elements remain.
"""

from __future__ import annotations

import os
from pathlib import Path

from pptx import Presentation


PROJECT_ROOT = Path(__file__).resolve().parents[1]
DEFAULT_TEMPLATE = Path(r"C:\Users\NSVLJ\Downloads\sampleppt (3).pptx")
DEFAULT_OUTPUT = PROJECT_ROOT / "Drought_Severity_Prediction_Presentation.pptx"


SLIDES = [
    {
        "title": "Predicting Drought Severity Using Climate and Vegetation Data",
        "body": [
            "Full-Stack Web Application + Machine Learning",
            "Flask REST API • React Dashboard • Neon PostgreSQL",
            "Inputs: Rainfall, Temperature, NDVI • Output: Severity (Low/Moderate/High)",
            "March 2026",
        ],
    },
    {
        "title": "Introduction & Problem Statement",
        "body": [
            "Drought impacts agriculture, water resources, and livelihoods.",
            "Traditional monitoring is reactive and slow.",
            "Goal: provide a fast, data-driven drought severity prediction tool.",
            "Users enter climate + vegetation indicators and get an instant severity level.",
        ],
    },
    {
        "title": "System Architecture",
        "body": [
            "Frontend: React (Login/Register + Dashboard + Charts)",
            "Backend: Flask REST API (Auth + Prediction)",
            "Database: Neon PostgreSQL (persistent users)",
            "ML: RandomForest model loaded from drought_model.pkl",
            "Flow: UI → API → Model → Prediction → UI visualization",
        ],
    },
    {
        "title": "Machine Learning Model",
        "body": [
            "Model: RandomForestClassifier",
            "Features: rainfall (mm), temperature (°C), NDVI (0–1)",
            "Target: severity class (Low, Moderate, High)",
            "Model persisted as: backend/drought_model.pkl",
            "API: POST /drought/predict returns severity label",
        ],
    },
    {
        "title": "Database Integration (Neon PostgreSQL)",
        "body": [
            "Users stored persistently in PostgreSQL (Neon)",
            "Passwords stored as hashed values (Werkzeug)",
            "ORM: Flask-SQLAlchemy",
            "On startup: db.create_all() creates required tables",
        ],
    },
    {
        "title": "Frontend Features (React)",
        "body": [
            "Pages: Login, Register, Dashboard",
            "Protected routing: dashboard requires login",
            "Prediction form inputs: rainfall, temperature, NDVI",
            "Visualization: Recharts bar chart for the input values",
            "Result card displays severity with color cues",
        ],
    },
    {
        "title": "Conclusion & Future Scope",
        "body": [
            "Delivered a complete end-to-end drought prediction web app.",
            "Improves usability with authentication, dashboard, and visualizations.",
            "Future: integrate real satellite/meteorological APIs for live data.",
            "Future: store prediction history per user for analytics.",
            "Future: containerize and deploy to cloud.",
        ],
    },
]


def _set_text(shape, text: str) -> None:
    if not shape.has_text_frame:
        return
    text_frame = shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = text


def _set_bullets(shape, lines: list[str]) -> None:
    if not shape.has_text_frame:
        return
    text_frame = shape.text_frame
    text_frame.clear()

    first = True
    for line in lines:
        if first:
            p = text_frame.paragraphs[0]
            first = False
        else:
            p = text_frame.add_paragraph()
        p.text = line
        p.level = 0


def _find_title_shape(slide):
    # Prefer placeholder titles, otherwise first text shape.
    for shape in slide.shapes:
        if shape.is_placeholder and shape.has_text_frame:
            try:
                phf = shape.placeholder_format
                if phf.type == 1:  # TITLE
                    return shape
            except Exception:
                continue
    for shape in slide.shapes:
        if shape.has_text_frame:
            return shape
    return None


def _find_body_shape(slide):
    # Prefer BODY placeholder, otherwise largest text box after title.
    candidates = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape.is_placeholder:
            try:
                if shape.placeholder_format.type in (2, 7):  # BODY or CONTENT
                    return shape
            except Exception:
                pass
        candidates.append(shape)

    # Heuristic: choose non-title text box with largest area.
    if not candidates:
        return None
    candidates_sorted = sorted(candidates, key=lambda s: (getattr(s, "width", 0) * getattr(s, "height", 0)), reverse=True)
    return candidates_sorted[0]


def generate(template_path: Path, output_path: Path) -> None:
    if not template_path.exists():
        raise FileNotFoundError(f"Template PPTX not found: {template_path}")

    prs = Presentation(str(template_path))

    # Ensure at least N slides.
    while len(prs.slides) < len(SLIDES):
        prs.slides.add_slide(prs.slide_layouts[1])

    # Update slide text for first N slides.
    for i, content in enumerate(SLIDES):
        slide = prs.slides[i]
        title_shape = _find_title_shape(slide)
        body_shape = _find_body_shape(slide)

        if title_shape:
            _set_text(title_shape, content["title"])

        if body_shape:
            # For title slide, keep as lines; for other slides, use bullets.
            if i == 0:
                _set_bullets(body_shape, content["body"])
            else:
                _set_bullets(body_shape, content["body"])

    # Optional: clear text from any remaining slides beyond our content.
    for i in range(len(SLIDES), len(prs.slides)):
        slide = prs.slides[i]
        for shape in slide.shapes:
            if shape.has_text_frame:
                shape.text_frame.clear()

    prs.save(str(output_path))


if __name__ == "__main__":
    template = Path(os.environ.get("PPT_TEMPLATE", str(DEFAULT_TEMPLATE)))
    output = Path(os.environ.get("PPT_OUTPUT", str(DEFAULT_OUTPUT)))
    generate(template, output)
    print(f"Generated PPT: {output}")
